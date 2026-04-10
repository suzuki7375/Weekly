import os
import re
import traceback
from collections import Counter
from copy import deepcopy
from pathlib import Path
import tkinter as tk
from tkinter import filedialog, messagebox

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt


def natural_sort_key(value: str):
    return [int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", value)]


def list_pptx_files(folder: Path):
    files = [p for p in folder.iterdir() if p.is_file() and p.suffix.lower() == ".pptx"]
    return sorted(files, key=lambda p: natural_sort_key(p.name))


def add_slide_copy(dest_prs: Presentation, source_slide):
    """將來源投影片內容複製到目標簡報的新投影片。"""
    blank_layout = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank_layout)

    for shape in source_slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue

        # 依關係類型決定 target：外部連結要用字串 URL，內部關係要用 target part。
        is_external = getattr(rel, "is_external", False)
        if is_external:
            target = rel.target_ref
        else:
            target = rel.target_part if hasattr(rel, "target_part") else rel._target

        # python-pptx 在不同版本的 Relationship API 名稱不同：
        # - 新版: add_relationship(...)
        # - 舊版: _add_relationship(...)
        rels = new_slide.part.rels
        if hasattr(rels, "add_relationship"):
            rels.add_relationship(rel.reltype, target, rel.rId, is_external=is_external)
        else:
            rels._add_relationship(rel.reltype, target, rel.rId, is_external=is_external)

    return new_slide


def extract_lines_from_ppt(ppt_path: Path):
    prs = Presentation(str(ppt_path))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                for line in shape.text.splitlines():
                    cleaned = line.strip(" \t•-\u2022")
                    if cleaned:
                        lines.append(cleaned)
    return lines


def summarize_person(lines, max_points=4):
    """簡易重點抽取：依長度與關鍵詞加權，選出最多 4 點。"""
    if not lines:
        return ["(未找到文字內容)"]

    unique = []
    seen = set()
    for line in lines:
        normalized = re.sub(r"\s+", " ", line)
        if normalized not in seen and len(normalized) >= 6:
            seen.add(normalized)
            unique.append(normalized)

    keywords = ["完成", "進行", "修正", "風險", "問題", "改善", "下週", "計畫", "支援", "上線"]
    score = Counter()
    for line in unique:
        base = min(len(line), 60)
        bonus = sum(8 for k in keywords if k in line)
        score[line] = base + bonus

    top = [item for item, _ in score.most_common(max_points)]
    return top if top else unique[:max_points]


def find_body_text_frame(slide):
    # 優先找內容 placeholder
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        phf = shape.placeholder_format
        if phf.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            if shape.has_text_frame:
                return shape.text_frame

    # 次選：第一個可寫入的文字框
    for shape in slide.shapes:
        if shape.has_text_frame:
            return shape.text_frame

    # 若都沒有，建立新的文字框
    box = slide.shapes.add_textbox(Pt(36), Pt(120), Pt(900), Pt(360))
    return box.text_frame


def write_summary_to_slide_2(template_prs: Presentation, summary_by_person):
    if len(template_prs.slides) < 2:
        raise ValueError("範例文件至少需要 2 頁，才能在第 2 頁寫入週報。")

    target_slide = template_prs.slides[1]

    # 設定標題
    if target_slide.shapes.title is not None:
        target_slide.shapes.title.text = "System team weekly Status"

    tf = find_body_text_frame(target_slide)
    tf.clear()

    for idx, (person, points) in enumerate(summary_by_person.items(), start=1):
        p = tf.paragraphs[0] if idx == 1 else tf.add_paragraph()
        p.text = f"{idx}. {person}"
        p.level = 0
        p.font.bold = True
        p.font.size = Pt(24)

        for point in points:
            sub = tf.add_paragraph()
            sub.text = f"• {point}"
            sub.level = 1
            sub.font.size = Pt(18)


def merge_and_write(template_file: Path, input_folder: Path, output_file: Path):
    if not template_file.exists():
        raise FileNotFoundError(f"找不到範例文件：{template_file}")
    if not input_folder.exists() or not input_folder.is_dir():
        raise NotADirectoryError(f"找不到整合資料資料夾：{input_folder}")

    input_files = list_pptx_files(input_folder)
    if not input_files:
        raise FileNotFoundError("整合資料資料夾內沒有 .pptx 檔案")

    template_prs = Presentation(str(template_file))

    summary_by_person = {}
    for ppt in input_files:
        person_name = ppt.stem
        lines = extract_lines_from_ppt(ppt)
        summary_by_person[person_name] = summarize_person(lines, max_points=4)

        source_prs = Presentation(str(ppt))
        for src_slide in source_prs.slides:
            add_slide_copy(template_prs, src_slide)

    write_summary_to_slide_2(template_prs, summary_by_person)
    template_prs.save(str(output_file))


class App:
    def __init__(self, root):
        self.root = root
        self.root.title("Weekly PPT 合併器")
        self.root.geometry("760x260")

        self.template_var = tk.StringVar()
        self.folder_var = tk.StringVar()
        self.output_var = tk.StringVar()

        self._build_ui()

    def _build_ui(self):
        pad_y = 14

        tk.Label(self.root, text="1. 範例文件", font=("Arial", 13, "bold")).grid(row=0, column=0, padx=10, pady=pad_y, sticky="w")
        tk.Entry(self.root, textvariable=self.template_var, width=68).grid(row=0, column=1, padx=6, pady=pad_y, sticky="we")
        tk.Button(self.root, text="選擇", command=self.pick_template).grid(row=0, column=2, padx=10, pady=pad_y)

        tk.Label(self.root, text="2. 整合資料(資料夾)", font=("Arial", 13, "bold")).grid(row=1, column=0, padx=10, pady=pad_y, sticky="w")
        tk.Entry(self.root, textvariable=self.folder_var, width=68).grid(row=1, column=1, padx=6, pady=pad_y, sticky="we")
        tk.Button(self.root, text="選擇", command=self.pick_folder).grid(row=1, column=2, padx=10, pady=pad_y)

        tk.Label(self.root, text="輸出文件", font=("Arial", 13, "bold")).grid(row=2, column=0, padx=10, pady=pad_y, sticky="w")
        tk.Entry(self.root, textvariable=self.output_var, width=68).grid(row=2, column=1, padx=6, pady=pad_y, sticky="we")
        tk.Button(self.root, text="另存", command=self.pick_output).grid(row=2, column=2, padx=10, pady=pad_y)

        tk.Button(
            self.root,
            text="3. 合併並撰寫",
            font=("Arial", 14, "bold"),
            bg="#0066cc",
            fg="white",
            command=self.on_merge_click,
            height=2,
        ).grid(row=3, column=0, columnspan=3, padx=10, pady=20, sticky="we")

        self.root.grid_columnconfigure(1, weight=1)

    def pick_template(self):
        path = filedialog.askopenfilename(filetypes=[("PowerPoint", "*.pptx")])
        if path:
            self.template_var.set(path)
            if not self.output_var.get().strip():
                default_out = str(Path(path).with_name(Path(path).stem + "_merged.pptx"))
                self.output_var.set(default_out)

    def pick_folder(self):
        path = filedialog.askdirectory()
        if path:
            self.folder_var.set(path)

    def pick_output(self):
        path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")])
        if path:
            self.output_var.set(path)

    def on_merge_click(self):
        template = Path(self.template_var.get().strip())
        folder = Path(self.folder_var.get().strip())
        output = Path(self.output_var.get().strip())

        if not template or not folder or not output:
            messagebox.showwarning("資料不足", "請先選擇範例文件、整合資料資料夾與輸出文件")
            return

        try:
            merge_and_write(template, folder, output)
            messagebox.showinfo("完成", f"已完成合併與撰寫：\n{output}")
        except Exception as exc:
            traceback.print_exc()
            messagebox.showerror("發生錯誤", f"{exc}")


if __name__ == "__main__":
    root = tk.Tk()
    app = App(root)
    root.mainloop()
